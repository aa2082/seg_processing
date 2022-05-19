from pptx import Presentation
from pptx.util import Inches


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]

path_to_data_folder = "/Users/ali/Desktop/exp1_seg_gif_with_roi/"
path_to_graphs = "/Users/ali/Desktop/mother_cell_segmentation_2/"
tr_num = 0
channels = ["mCherry","YFP"]
graph_num = 1
for n_fov in range(0,2):
    for tr_num in range(0,2):
        FOV = "xy{:03d}".format(n_fov)
        slide = prs.slides.add_slide(title_slide_layout)


        graph = path_to_graphs+str(graph_num)+".png"

        pic = slide.shapes.add_picture(graph, Inches(2), Inches(1))

        gif_mCherry = path_to_data_folder+FOV+"/Masks/stack_"+FOV+"_tr_"+str(tr_num)+"_mCherry.gif"
        gif_YFP = path_to_data_folder+FOV+"/Masks/stack_"+FOV+"_tr_"+str(tr_num)+"_YFP.gif"

        pic = slide.shapes.add_picture(gif_mCherry, Inches(1), Inches(1),height=Inches(5))
        pic = slide.shapes.add_picture(gif_YFP, Inches(1.5), Inches(1),height=Inches(5))

        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        graph_num+=1

prs.save('test.pptx')